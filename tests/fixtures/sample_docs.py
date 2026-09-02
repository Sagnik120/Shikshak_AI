"""Sample document fixtures generator for comprehensive RAG testing."""

import io
from typing import Dict, Any


def get_physics_notes_markdown() -> str:
    """Multi-section markdown text with clear hierarchy."""
    return """# Electromagnetism and Circuits

## Section 1: Electric Current and Voltage
Electric current is defined as the rate of flow of electric charge through a conductor.
The SI unit of current is the Ampere (A).
Voltage, also known as electric potential difference, is the work done per unit charge in moving a charge between two points.
The formula for electric potential is V = W / Q.

## Section 2: Ohm's Law and Resistance
Ohm's law states that the current flowing through a conductor is directly proportional to the potential difference across its ends, provided the physical conditions (such as temperature) remain constant.
Mathematically, V = I * R, where R is resistance measured in Ohms (Ω).
Resistance depends on the material resistivity, length of the conductor, and cross-sectional area: R = ρ * (L / A).

## Section 3: Electrical Power and Joule's Heating
Electric power is the rate at which electrical energy is transferred by an electric circuit.
Power P = V * I = I^2 * R = V^2 / R.
Joule's law of heating states that heat produced H = I^2 * R * t.
"""


def get_hindi_biology_markdown() -> str:
    """Hindi language biology document to test Devanagari language detection and multilingual RAG."""
    return """# प्रकाश संश्लेषण (Photosynthesis)

## खंड 1: परिचय
प्रकाश संश्लेषण वह जैविक प्रक्रिया है जिसके द्वारा हरे पौधे सूर्य के प्रकाश, जल और कार्बन डाइऑक्साइड का उपयोग करके अपना भोजन (ग्लूकोज) बनाते हैं।

## खंड 2: रासायनिक समीकरण
इस प्रक्रिया का मुख्य रासायनिक समीकरण है:
6CO2 + 6H2O + प्रकाश ऊर्जा -> C6H12O6 + 6O2
इस प्रक्रिया में क्लोरोफिल एक महत्वपूर्ण वर्णक के रूप में कार्य करता है।
"""


def get_single_paragraph_doc() -> str:
    """Short single-paragraph document for edge-case testing."""
    return "Mitochondria are membrane-bound cell organelles that generate most of the chemical energy needed to power the cell's biochemical reactions."


def get_unstructured_wall_of_text() -> str:
    """Document without any headers or bullet points."""
    return (
        "Classical mechanics describes the motion of macroscopic objects from projectiles to parts of machinery. "
        "It is one of the oldest subjects in science. "
        "Isaac Newton formulated three laws of motion in 1687. "
        "The first law is the law of inertia. The second law defines force as mass times acceleration. "
        "The third law states that for every action there is an equal and opposite reaction."
    )


def create_sample_docx_bytes() -> bytes:
    """Create a minimal valid DOCX document in-memory using python-docx."""
    try:
        from docx import Document
        doc = Document()
        doc.add_heading("Quantum Mechanics Overview", level=1)
        doc.add_paragraph("Quantum mechanics is the study of matter and radiation at an atomic level.")
        doc.add_heading("Wave-Particle Duality", level=2)
        doc.add_paragraph("Light exhibits both wave and particle characteristics, described by de Broglie wavelength λ = h / p.")
        
        table = doc.add_table(rows=2, cols=2)
        table.rows[0].cells[0].text = "Particle"
        table.rows[0].cells[1].text = "Mass (kg)"
        table.rows[1].cells[0].text = "Electron"
        table.rows[1].cells[1].text = "9.109e-31"

        buf = io.BytesIO()
        doc.save(buf)
        return buf.getvalue()
    except Exception:
        # If python-docx not installed yet, return raw simulated string
        return b"Quantum Mechanics Overview\nWave-Particle Duality"


def create_sample_pptx_bytes() -> bytes:
    """Create a minimal valid PPTX presentation in-memory using python-pptx."""
    try:
        from pptx import Presentation
        from pptx.util import Inches
        prs = Presentation()
        slide_layout = prs.slide_layouts[0] # title slide
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        title.text = "Thermodynamics Lecture"
        subtitle.text = "First and Second Laws of Thermodynamics"

        slide_layout_2 = prs.slide_layouts[1] # bullet slide
        slide2 = prs.slides.add_slide(slide_layout_2)
        slide2.shapes.title.text = "Zeroth and First Law"
        tf = slide2.shapes.placeholders[1].text_frame
        tf.text = "Zeroth law establishes thermal equilibrium and temperature."
        p = tf.add_paragraph()
        p.text = "First law states energy conservation: ΔU = Q - W."

        buf = io.BytesIO()
        prs.save(buf)
        return buf.getvalue()
    except Exception:
        return b"Thermodynamics Lecture\nZeroth and First Law"
