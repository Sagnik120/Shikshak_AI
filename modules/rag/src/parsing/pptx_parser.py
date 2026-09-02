"""PPTX presentation parser extracting text per slide, slide titles, and table structures."""

from __future__ import annotations

import io
import logging
from typing import List, Tuple

from modules.rag.src.models import RawSection

logger = logging.getLogger(__name__)


def parse_pptx(file_bytes: bytes) -> Tuple[List[RawSection], List[str]]:
    """Parse a PPTX file into raw sections per slide, capturing slide titles and slide numbers.

    Args:
        file_bytes: Raw bytes of the uploaded PPTX presentation.

    Returns:
        Tuple of (raw_sections, detected_chapters)
    """
    sections: List[RawSection] = []
    chapters: List[str] = []

    try:
        from pptx import Presentation
        prs = Presentation(io.BytesIO(file_bytes))

        for idx, slide in enumerate(prs.slides):
            slide_number = idx + 1
            slide_title: str | None = None
            slide_text_blocks: List[str] = []

            # Try to get slide title placeholder
            try:
                if slide.shapes.title and slide.shapes.title.text.strip():
                    slide_title = slide.shapes.title.text.strip()
                    if slide_title and slide_title not in chapters:
                        chapters.append(slide_title)
            except Exception:
                pass

            # Iterate through all shapes in the slide
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text.strip()
                        # If title wasn't found from shape.title, use first prominent text
                        if not slide_title and text and len(text) < 60:
                            slide_title = text
                            if slide_title not in chapters:
                                chapters.append(slide_title)
                        elif text and text != slide_title:
                            slide_text_blocks.append(text)

                elif shape.has_table:
                    for row in shape.table.rows:
                        row_cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                        if row_cells:
                            slide_text_blocks.append(" | ".join(row_cells))

            full_slide_text = ""
            if slide_title:
                full_slide_text = f"Title: {slide_title}\n" + "\n".join(slide_text_blocks)
            else:
                full_slide_text = "\n".join(slide_text_blocks)

            if full_slide_text.strip():
                sections.append(
                    RawSection(
                        section_title=slide_title,
                        page_or_slide=slide_number,
                        raw_text=full_slide_text.strip(),
                        metadata={"slide_number": slide_number}
                    )
                )

    except Exception as e:
        logger.error(f"Failed to parse PPTX: {e}")

    return sections, chapters
